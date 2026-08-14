import os
import uuid
import tempfile
from typing import Optional
from fastapi import UploadFile
from langchain_chroma import Chroma
from langchain_google_genai import GoogleGenerativeAIEmbeddings
from langchain_text_splitters import RecursiveCharacterTextSplitter
from langchain_community.document_loaders import PyMuPDFLoader, Docx2txtLoader, TextLoader
from app.core.config import settings
from app.core.prompt_security import sanitize_rag_context

# Initialize Embeddings
embeddings = GoogleGenerativeAIEmbeddings(
    model="models/gemini-embedding-2",
    google_api_key=settings.GEMINI_API_KEY,
)

# Initialize ChromaDB persistent client
CHROMA_PERSIST_DIR = os.path.join(os.getcwd(), "chroma_db")
vectorstore = Chroma(
    collection_name="student_documents",
    embedding_function=embeddings,
    persist_directory=CHROMA_PERSIST_DIR,
)

cache_vectorstore = Chroma(
    collection_name="semantic_cache",
    embedding_function=embeddings,
    persist_directory=CHROMA_PERSIST_DIR,
)


async def process_and_store_document(
    file: UploadFile,
    document_id: str,
    user_id: str,
    raw_content: Optional[bytes] = None,
):
    """Process an uploaded file, chunk it, and store vectors in ChromaDB."""

    file_ext = os.path.splitext(file.filename or "")[1].lower()

    # Use pre-read bytes if provided (avoids consuming the stream twice)
    if raw_content is None:
        raw_content = await file.read()

    temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=file_ext)
    pptx_tmp = None
    try:
        temp_file.write(raw_content)
        temp_file.close()

        # ── Select loader ──────────────────────────────────────────────────────
        loader = None
        if file_ext == ".pdf":
            loader = PyMuPDFLoader(temp_file.name)
        elif file_ext in (".docx", ".doc"):
            loader = Docx2txtLoader(temp_file.name)
        elif file_ext == ".txt":
            loader = TextLoader(temp_file.name, encoding="utf-8")
        elif file_ext in (".pptx", ".ppt"):
            # Extract text from PowerPoint using python-pptx
            from pptx import Presentation

            prs = Presentation(temp_file.name)
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_runs.append(shape.text)
            text = "\n".join(text_runs)

            pptx_tmp = tempfile.NamedTemporaryFile(delete=False, suffix=".txt")
            pptx_tmp.write(text.encode("utf-8"))
            pptx_tmp.close()
            loader = TextLoader(pptx_tmp.name, encoding="utf-8")
        else:
            raise ValueError(f"Unsupported file format: {file_ext}")

        docs = loader.load()

        # ── Split into chunks ──────────────────────────────────────────────────
        text_splitter = RecursiveCharacterTextSplitter(
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len,
        )
        chunks = text_splitter.split_documents(docs)

        if not chunks:
            return  # Empty document — nothing to store

        # ── Attach metadata & sanitize content ──────────────────────────────
        for chunk in chunks:
            chunk.page_content = sanitize_rag_context(chunk.page_content)
            chunk.metadata.update({
                "document_id": document_id,
                "user_id": user_id,
                "source": file.filename,
            })

        # ── Store in ChromaDB ──────────────────────────────────────────────────
        vectorstore.add_documents(chunks)

    finally:
        # Cleanup temp files
        if os.path.exists(temp_file.name):
            os.remove(temp_file.name)
        if pptx_tmp and os.path.exists(pptx_tmp.name):
            os.remove(pptx_tmp.name)


def search_documents(
    query: str,
    user_id: str,
    document_id: str = None,
    top_k: int = 3,
):
    """Retrieve relevant chunks for a specific user and optionally a specific document."""
    if document_id:
        filter_dict = {
            "$and": [
                {"user_id": user_id},
                {"document_id": document_id}
            ]
        }
    else:
        filter_dict = {"user_id": user_id}

    # Fetch more candidates to allow for deduplication and filtering
    results_with_scores = vectorstore.similarity_search_with_score(query, k=top_k * 2, filter=filter_dict)
    
    unique_contents = set()
    final_results = []
    
    for doc, score in results_with_scores:
        # Filter low relevance: ChromaDB default is L2 distance, lower is better.
        if score < 1.2:
            content_hash = hash(doc.page_content)
            if content_hash not in unique_contents:
                unique_contents.add(content_hash)
                doc.page_content = sanitize_rag_context(doc.page_content)
                final_results.append(doc)
                if len(final_results) >= top_k:
                    break

    return final_results


def check_cache(query: str, user_id: str, threshold: float = 0.2):
    """Check semantic cache for highly similar queries."""
    filter_dict = {"user_id": user_id}
    results = cache_vectorstore.similarity_search_with_score(query, k=1, filter=filter_dict)
    if results:
        doc, score = results[0]
        # Very strict threshold to ensure questions mean the same thing
        if score < threshold:
            return doc.metadata.get("response")
    return None


def set_cache(query: str, response: str, user_id: str):
    """Store the query and its response in the semantic cache."""
    cache_vectorstore.add_texts(
        texts=[query],
        metadatas=[{"response": response, "user_id": user_id}]
    )
